"""Read PowerPoint action handler."""

from __future__ import annotations

from pathlib import Path

from tools.file_ops.helpers import _safe_resolve
from tools.file_ops._registry import register_action


@register_action(
    "file",
    "read_pptx",
    help_text="""Read a PowerPoint file using python-pptx.
Required: path
Optional: max_chars (default 50000)
Returns: {text, slides, slide_count, truncated}""",
    examples=[
        'file(action="read_pptx", path="presentation.pptx")',
    ],
)
def _handle_read_pptx(path: str = "", max_chars: int = 50_000, trace_id: str = "", **kwargs) -> dict:
    """Read a PowerPoint file using python-pptx."""
    p, err = _safe_resolve(path)
    if err:
        return {"status": "error", "error": err}
    if not p.exists():
        return {"status": "error", "error": f"File not found: {p}"}
    if p.suffix.lower() != ".pptx":
        return {"status": "error", "error": f"Not a .pptx file: {p.name}"}

    try:
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation(str(p))
        slides = []

        for i, slide in enumerate(prs.slides, 1):
            texts = []
            images = 0
            tables = 0

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
                if hasattr(shape, "table"):
                    tables += 1
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        texts.append(" | ".join(cells))
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    images += 1

            layout = ""
            try:
                layout = slide.slide_layout.name
            except Exception:
                pass

            slides.append({
                "slide": i,
                "layout": layout,
                "texts": texts,
                "images": images,
                "tables": tables,
            })

        flat = "\n\n".join(
            f"--- Slide {s['slide']} ---\n" + "\n".join(s["texts"])
            for s in slides
            if s["texts"]
        )

        truncated = len(flat) > max_chars
        if truncated:
            flat = flat[:max_chars] + "\n\n[...truncated]"

        return {
            "status": "success",
            "path": str(p),
            "text": flat,
            "slides": slides,
            "slide_count": len(slides),
            "truncated": truncated,
        }
    except ImportError:
        return {"status": "error",
                "error": "python-pptx not installed. Run: pip install python-pptx"}
    except Exception as e:
        return {"status": "error", "error": f"PPTX read failed: {type(e).__name__}: {e}"}
