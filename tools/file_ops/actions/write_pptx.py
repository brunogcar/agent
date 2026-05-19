"""
Write PowerPoint action handler.
"""

from __future__ import annotations

import json as _json
from pathlib import Path

from core.config import cfg
from tools.file_ops.helpers import _safe_resolve

def _handle_write_pptx(path: str = "", content: str = "", **kwargs) -> dict:
    """Write data to a PowerPoint file using python-pptx."""
    p, err = _safe_resolve(path)
    if err:
        return {"status": "error", "error": err}
    if not content:
        return {"status": "error", "error": "content is required for write_pptx"}
    if p.suffix.lower() != ".pptx":
        p = p.with_suffix(".pptx")

    # Parse content if it's a string
    if isinstance(content, str):
        try:
            content = _json.loads(content)
        except Exception:
            return {"status": "error",
                    "error": "content string could not be parsed as JSON"}

    if not isinstance(content, list):
        return {"status": "error",
                "error": "content must be a list of slide dicts: [{title, body/bullets}, ...]"}

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Color scheme
        DARK = RGBColor(0x2C, 0x3E, 0x50)
        ACCENT = RGBColor(0x34, 0x98, 0xDB)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        LIGHT = RGBColor(0xEC, 0xF0, 0xF1)

        layouts = {n.name: n for n in prs.slide_layouts}

        def _add_rect(slide, left, top, width, height, color):
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                Inches(left), Inches(top), Inches(width), Inches(height),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()
            return shape

        def _add_text(slide, text, left, top, width, height,
                      font_size=18, bold=False, color=None, align="left"):
            txb = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height))
            tf = txb.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.alignment = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
            }.get(align, PP_ALIGN.LEFT)
            run = para.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            if color:
                run.font.color.rgb = color
            return txb

        # Build slides
        for i, spec in enumerate(content):
            slide_title = spec.get("title", f"Slide {i+1}")
            slide_body = spec.get("body", "")
            slide_bullets = spec.get("bullets", [])
            slide_layout = spec.get("layout", "content").lower()
            slide_notes = spec.get("notes", "")
            is_title_slide = (slide_layout == "title" or i == 0 and slide_layout != "content")

            # Use blank layout for full control
            blank = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(blank)

            if is_title_slide and i == 0:
                # Title slide - dark background, centered
                _add_rect(slide, 0, 0, 13.33, 7.5, DARK)
                _add_rect(slide, 0, 3.0, 13.33, 0.06, ACCENT)
                _add_text(slide, slide_title,
                          1, 2.2, 11.33, 1.4,
                          font_size=40, bold=True, color=WHITE, align="center")
                if slide_body:
                    _add_text(slide, slide_body,
                              1, 3.8, 11.33, 0.8,
                              font_size=20, color=LIGHT, align="center")
                if spec.get("subtitle"):
                    _add_text(slide, spec["subtitle"],
                              1, 4.6, 11.33, 0.7,
                              font_size=16, color=ACCENT, align="center")
            else:
                # Content slide
                _add_rect(slide, 0, 0, 13.33, 1.1, DARK)
                _add_rect(slide, 0, 1.1, 13.33, 0.05, ACCENT)

                _add_text(slide, slide_title,
                          0.3, 0.15, 12.73, 0.85,
                          font_size=26, bold=True, color=WHITE)
                _add_text(slide, str(i + 1),
                          12.5, 0.25, 0.6, 0.5,
                          font_size=13, color=ACCENT, align="right")

                content_top = 1.35

                if slide_bullets:
                    txb = slide.shapes.add_textbox(
                        Inches(0.5), Inches(content_top),
                        Inches(12.33), Inches(5.8))
                    tf = txb.text_frame
                    tf.word_wrap = True

                    for j, bullet in enumerate(slide_bullets):
                        if j == 0:
                            para = tf.paragraphs[0]
                        else:
                            para = tf.add_paragraph()

                        if isinstance(bullet, dict):
                            bullet_text = bullet.get("text", "")
                            bullet_level = bullet.get("level", 0)
                        else:
                            bullet_text = str(bullet)
                            bullet_level = 0

                        para.level = bullet_level
                        run = para.add_run()
                        run.text = bullet_text
                        run.font.size = Pt(18 - bullet_level * 2)
                        run.font.color.rgb = DARK
                        para.space_before = Pt(6)

                elif slide_body:
                    _add_text(slide, slide_body,
                              0.5, content_top, 12.33, 5.8,
                              font_size=18, color=DARK)

            # Speaker notes
            if slide_notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_notes

        p.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(p))

        return {
            "status": "success",
            "path": str(p),
            "size": p.stat().st_size,
            "slide_count": len(content),
        }
    except ImportError:
        return {"status": "error",
                "error": "python-pptx not installed. Run: pip install python-pptx"}
    except Exception as e:
        return {"status": "error", "error": f"PPTX write failed: {type(e).__name__}: {e}"}