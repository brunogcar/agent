"""
tools/file_ops.py — File meta-tool.

Replaces: fs MCP server for most operations (fs still exists for safety isolation)
The LLM sees ONE tool: file(action, ...)

Actions:
  read      → read a single file
  write     → write content to a file (auto-backup)
  list      → list directory contents
  backup    → copy file with .bak suffix
  read_many → read multiple files concurrently
  search    → full-text search across workspace (SQLite FTS)
  read_pdf  → extract text from PDF using pdfplumber

All paths are resolved relative to workspace_root unless absolute.
Paths outside workspace_root and agent_root are rejected for safety.
"""

from __future__ import annotations

import sqlite3
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import Optional

from core.config import cfg
from registry import tool

# ── Path safety ───────────────────────────────────────────────────────────────

_ALLOWED_ROOTS = None

def _allowed_roots() -> list[Path]:
    global _ALLOWED_ROOTS
    if _ALLOWED_ROOTS is None:
        _ALLOWED_ROOTS = [
            cfg.workspace_root.resolve(),
            cfg.agent_root.resolve(),
        ]
    return _ALLOWED_ROOTS


def _resolve(path_str: str) -> Optional[Path]:
    """
    Resolve a path safely.
    - Absolute paths are used as-is (if within allowed roots)
    - Relative paths are resolved from workspace_root
    Returns None if the path escapes allowed roots.
    """
    p = Path(path_str)
    if not p.is_absolute():
        p = cfg.workspace_root / p

    resolved = p.resolve()
    for root in _allowed_roots():
        try:
            resolved.relative_to(root)
            return resolved
        except ValueError:
            continue
    return None  # path escapes allowed roots


def _safe_resolve(path_str: str) -> tuple[Optional[Path], str]:
    """Returns (resolved_path, error_message). error is "" on success."""
    if not path_str:
        return None, "path is required"
    p = _resolve(path_str)
    if p is None:
        return None, (
            f"Path '{path_str}' is outside allowed directories. "
            f"Use paths within workspace ({cfg.workspace_root}) or agent ({cfg.agent_root})."
        )
    return p, ""


# ── SQLite FTS index ──────────────────────────────────────────────────────────

_INDEX_DB: Optional[sqlite3.Connection] = None

def _get_index() -> sqlite3.Connection:
    global _INDEX_DB
    if _INDEX_DB is None:
        cfg.workspace_index.mkdir(parents=True, exist_ok=True)
        db_path = cfg.workspace_index / "fts.db"
        _INDEX_DB = sqlite3.connect(str(db_path), check_same_thread=False)
        _INDEX_DB.execute("""
            CREATE VIRTUAL TABLE IF NOT EXISTS files_fts
            USING fts5(path, content, tokenize='porter ascii')
        """)
        _INDEX_DB.execute("""
            CREATE TABLE IF NOT EXISTS files_meta (
                path TEXT PRIMARY KEY,
                mtime REAL,
                size  INTEGER
            )
        """)
        _INDEX_DB.commit()
    return _INDEX_DB


def _index_file(path: Path) -> bool:
    """Add or update a file in the FTS index. Returns True on success."""
    try:
        stat  = path.stat()
        mtime = stat.st_mtime
        size  = stat.st_size

        if size > 500_000:  # skip files > 500KB
            return False

        db = _get_index()

        # Check if up to date
        row = db.execute(
            "SELECT mtime FROM files_meta WHERE path = ?", (str(path),)
        ).fetchone()
        if row and abs(row[0] - mtime) < 0.01:
            return True  # already indexed and unchanged

        text = path.read_text(encoding="utf-8", errors="replace")

        db.execute("DELETE FROM files_fts WHERE path = ?", (str(path),))
        db.execute("INSERT INTO files_fts(path, content) VALUES (?, ?)", (str(path), text))
        db.execute(
            "INSERT OR REPLACE INTO files_meta(path, mtime, size) VALUES (?, ?, ?)",
            (str(path), mtime, size),
        )
        db.commit()
        return True
    except Exception:
        return False


def _build_index(root: Path, extensions: set[str] = None) -> int:
    """Index all text files under root. Returns count indexed."""
    if extensions is None:
        extensions = {".py", ".md", ".txt", ".json", ".yaml", ".yml", ".toml", ".cfg", ".ini"}

    count = 0
    for p in root.rglob("*"):
        if p.is_file() and p.suffix in extensions:
            # Skip hidden, cache, git
            parts = p.parts
            if any(part.startswith(".") or part == "__pycache__" for part in parts):
                continue
            if _index_file(p):
                count += 1
    return count


# ── Read helpers ──────────────────────────────────────────────────────────────

def _read_file(path: Path, max_chars: int = 50_000) -> dict:
    """Read a single file and return structured result."""
    if not path.exists():
        return {"status": "error", "error": f"File not found: {path}"}
    if not path.is_file():
        return {"status": "error", "error": f"Not a file: {path}"}

    stat = path.stat()
    if stat.st_size == 0:
        return {
            "status": "success", "path": str(path),
            "content": "", "size": 0, "lines": 0, "truncated": False,
        }

    try:
        text = path.read_text(encoding="utf-8", errors="replace")
        truncated = len(text) > max_chars
        if truncated:
            text = text[:max_chars] + f"\n\n[...truncated — {stat.st_size} bytes total]"

        return {
            "status":    "success",
            "path":      str(path),
            "content":   text,
            "size":      stat.st_size,
            "lines":     text.count("\n") + 1,
            "truncated": truncated,
            "extension": path.suffix,
        }
    except Exception as e:
        return {"status": "error", "error": f"Read failed: {e}", "path": str(path)}


# ── Meta-tool ─────────────────────────────────────────────────────────────────

@tool
def file(
    action:    str,
    path:      str        = "",
    content:   str        = "",
    paths:     list       = None,
    mode:      str        = "full",
    query:     str        = "",
    max_chars: int        = 50_000,
    max_results: int      = 10,
    title:       str      = "",
) -> dict:
    """
    File tool — read, write, search, and manage files.

    action: "read" | "write" | "list" | "backup" | "read_many" | "search" |
            "read_pdf" | "write_pdf" |
            "read_docx" | "write_docx" |
            "read_xlsx" | "write_xlsx" |
            "read_pptx" | "write_pptx"

    read
        Read a single file. Paths relative to workspace root.
        Required: path
        Optional: max_chars (default 50000)
        Returns:  {content, size, lines, truncated}

    write
        Write content to a file. Auto-creates parent directories.
        If file exists, a .bak backup is created automatically.
        Required: path, content
        Returns:  {path, size, backup_path}

    list
        List contents of a directory.
        Required: path (directory)
        Returns:  {entries: [{name, type, size, modified}]}

    backup
        Copy a file with .bak suffix (manual backup).
        Required: path
        Returns:  {original, backup}

    read_many
        Read multiple files concurrently. Returns all contents in one call.
        mode: "full" (default) | "summary" (first 500 chars per file)
        Required: paths (list of path strings)
        Optional: mode, max_chars
        Returns:  {files: [{path, content, size, error}], count}

    search
        Full-text search across workspace files.
        Builds/updates the index automatically on first use.
        Required: query
        Optional: max_results (default 10)
        Returns:  {results: [{path, snippet, rank}], count}

    read_pdf
        Extract text from a PDF file using pdfplumber.
        Required: path
        Optional: max_chars
        Returns:  {text, pages, truncated}

    Examples:
        file(action="read", path="scripts/analysis.py")
        file(action="write", path="output/report.md", content="# Report\\n...")
        file(action="list", path=".")
        file(action="read_many", paths=["a.py", "b.py", "c.py"], mode="summary")
        file(action="search", query="ChromaDB collection", max_results=5)
        file(action="read_pdf", path="docs/manual.pdf")
    """
    action = action.strip().lower()

    # ── read ──────────────────────────────────────────────────────────────────
    if action == "read":
        p, err = _safe_resolve(path)
        if err:
            return {"status": "error", "error": err}
        return _read_file(p, max_chars)

    # ── write ─────────────────────────────────────────────────────────────────
    if action == "write":
        p, err = _safe_resolve(path)
        if err:
            return {"status": "error", "error": err}
        if not content and content != "":
            return {"status": "error", "error": "content is required for write"}
        if cfg.is_protected(p):
            return {"status": "error", "error": f"'{p.name}' is protected — edit manually"}

        backup_path = ""
        p.parent.mkdir(parents=True, exist_ok=True)

        # Auto-backup if file exists
        if p.exists():
            bak = p.with_suffix(p.suffix + ".bak")
            import shutil
            shutil.copy2(p, bak)
            backup_path = str(bak)

        try:
            p.write_text(content, encoding="utf-8")
            return {
                "status":      "success",
                "path":        str(p),
                "size":        len(content.encode("utf-8")),
                "backup_path": backup_path,
                "lines":       content.count("\n") + 1,
            }
        except Exception as e:
            return {"status": "error", "error": f"Write failed: {e}"}

    # ── list ──────────────────────────────────────────────────────────────────
    if action == "list":
        p, err = _safe_resolve(path or ".")
        if err:
            return {"status": "error", "error": err}
        if not p.is_dir():
            return {"status": "error", "error": f"Not a directory: {p}"}

        entries = []
        try:
            for item in sorted(p.iterdir()):
                stat = item.stat()
                entries.append({
                    "name":     item.name,
                    "type":     "dir" if item.is_dir() else "file",
                    "size":     stat.st_size if item.is_file() else 0,
                    "modified": time.strftime(
                        "%Y-%m-%d %H:%M", time.localtime(stat.st_mtime)
                    ),
                    "extension": item.suffix if item.is_file() else "",
                })
        except PermissionError as e:
            return {"status": "error", "error": f"Permission denied: {e}"}

        return {
            "status":  "success",
            "path":    str(p),
            "entries": entries,
            "count":   len(entries),
        }

    # ── backup ────────────────────────────────────────────────────────────────
    if action == "backup":
        import shutil
        p, err = _safe_resolve(path)
        if err:
            return {"status": "error", "error": err}
        if not p.exists():
            return {"status": "error", "error": f"File not found: {p}"}

        ts  = time.strftime("%Y%m%d_%H%M%S")
        bak = p.with_name(f"{p.stem}_{ts}{p.suffix}.bak")
        try:
            shutil.copy2(p, bak)
            return {"status": "success", "original": str(p), "backup": str(bak)}
        except Exception as e:
            return {"status": "error", "error": f"Backup failed: {e}"}

    # ── read_many ─────────────────────────────────────────────────────────────
    if action == "read_many":
        if not paths:
            return {"status": "error", "error": "paths list is required for read_many"}

        summary_chars = 500 if mode == "summary" else max_chars

        def _read_one(path_str: str) -> dict:
            p, err = _safe_resolve(path_str)
            if err:
                return {"path": path_str, "error": err, "content": "", "size": 0}
            result = _read_file(p, summary_chars)
            return {
                "path":    path_str,
                "content": result.get("content", ""),
                "size":    result.get("size", 0),
                "lines":   result.get("lines", 0),
                "error":   result.get("error", ""),
            }

        results = []
        # ThreadPoolExecutor for concurrent reads
        with ThreadPoolExecutor(max_workers=min(len(paths), 8)) as executor:
            futures = {executor.submit(_read_one, p): p for p in paths}
            for future in as_completed(futures):
                results.append(future.result())

        # Restore original order
        order  = {p: i for i, p in enumerate(paths)}
        results.sort(key=lambda r: order.get(r["path"], 999))

        total_size = sum(r["size"] for r in results)
        errors     = [r["path"] for r in results if r["error"]]

        return {
            "status":     "success",
            "files":      results,
            "count":      len(results),
            "total_size": total_size,
            "errors":     errors,
            "mode":       mode,
        }

    # ── search ────────────────────────────────────────────────────────────────
    if action == "search":
        if not query:
            return {"status": "error", "error": "query is required for search"}

        # Build/refresh index
        indexed = _build_index(cfg.workspace_root)

        try:
            db = _get_index()
            rows = db.execute(
                """
                SELECT path, snippet(files_fts, 1, '[', ']', '...', 20) as snippet, rank
                FROM files_fts
                WHERE content MATCH ?
                ORDER BY rank
                LIMIT ?
                """,
                (query, max_results),
            ).fetchall()

            results = [
                {"path": row[0], "snippet": row[1], "rank": round(abs(row[2]), 4)}
                for row in rows
            ]

            return {
                "status":        "success",
                "query":         query,
                "results":       results,
                "count":         len(results),
                "indexed_files": indexed,
            }
        except Exception as e:
            return {"status": "error", "error": f"Search failed: {e}"}

    # ── read_pdf ──────────────────────────────────────────────────────────────
    if action == "read_pdf":
        p, err = _safe_resolve(path)
        if err:
            return {"status": "error", "error": err}
        if not p.exists():
            return {"status": "error", "error": f"File not found: {p}"}
        if p.suffix.lower() != ".pdf":
            return {"status": "error", "error": f"Not a PDF file: {p.name}"}

        try:
            import pdfplumber

            pages_text = []
            with pdfplumber.open(str(p)) as pdf:
                total_pages = len(pdf.pages)
                for page in pdf.pages:
                    text = page.extract_text() or ""
                    if text.strip():
                        pages_text.append(text)

            full_text = "\n\n".join(pages_text)
            truncated = len(full_text) > max_chars
            if truncated:
                full_text = full_text[:max_chars] + f"\n\n[...truncated — {total_pages} pages total]"

            return {
                "status":    "success",
                "path":      str(p),
                "text":      full_text,
                "pages":     total_pages,
                "truncated": truncated,
                "word_count": len(full_text.split()),
            }
        except ImportError:
            return {"status": "error", "error": "pdfplumber not installed. Run: pip install pdfplumber"}
        except Exception as e:
            return {"status": "error", "error": f"PDF read failed: {type(e).__name__}: {e}"}

    # ── read_docx ─────────────────────────────────────────────────────────────
    if action == "read_docx":
        p, err = _safe_resolve(path)
        if err:
            return {"status": "error", "error": err}
        if not p.exists():
            return {"status": "error", "error": f"File not found: {p}"}
        if p.suffix.lower() != ".docx":
            return {"status": "error", "error": f"Not a .docx file: {p.name}"}
        try:
            from docx import Document
            from docx.oxml.ns import qn

            doc      = Document(str(p))
            sections = []
            tables   = []

            for elem in doc.element.body:
                tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag

                if tag == "p":
                    # Paragraph — detect heading style
                    from docx.text.paragraph import Paragraph
                    para  = Paragraph(elem, doc)
                    text  = para.text.strip()
                    style = para.style.name if para.style else ""
                    if not text:
                        continue
                    if style.startswith("Heading"):
                        level = style.replace("Heading ", "").strip()
                        sections.append({"type": "heading", "level": level, "text": text})
                    else:
                        sections.append({"type": "paragraph", "text": text})

                elif tag == "tbl":
                    from docx.table import Table
                    tbl  = Table(elem, doc)
                    rows = []
                    for row in tbl.rows:
                        rows.append([c.text.strip() for c in row.cells])
                    tables.append(rows)
                    sections.append({"type": "table", "rows": rows})

            # Flat text version for easy LLM consumption
            flat = "\n".join(
                ("#" * int(s.get("level", 1)) + " " + s["text"])
                if s["type"] == "heading"
                else s["text"]
                if s["type"] == "paragraph"
                else "[TABLE: " + " | ".join(s["rows"][0]) + " ...]"
                if s["type"] == "table" and s["rows"]
                else ""
                for s in sections
            ).strip()

            truncated = len(flat) > max_chars
            if truncated:
                flat = flat[:max_chars] + f"\n\n[...truncated]"

            return {
                "status":     "success",
                "path":       str(p),
                "text":       flat,
                "sections":   sections,
                "tables":     len(tables),
                "paragraphs": len([s for s in sections if s["type"] == "paragraph"]),
                "truncated":  truncated,
            }
        except ImportError:
            return {"status": "error", "error": "python-docx not installed. Run: pip install python-docx"}
        except Exception as e:
            return {"status": "error", "error": f"DOCX read failed: {type(e).__name__}: {e}"}

    # ── read_xlsx ─────────────────────────────────────────────────────────────
    if action == "read_xlsx":
        p, err = _safe_resolve(path)
        if err:
            return {"status": "error", "error": err}
        if not p.exists():
            return {"status": "error", "error": f"File not found: {p}"}
        if p.suffix.lower() not in (".xlsx", ".xls", ".xlsm"):
            return {"status": "error", "error": f"Not an Excel file: {p.name}"}
        try:
            import pandas as pd

            xl       = pd.ExcelFile(str(p))
            sheets   = xl.sheet_names
            result   = {}

            for sheet in sheets:
                df = xl.parse(sheet)
                # Cap rows to avoid token explosion
                MAX_ROWS = 200
                truncated_sheet = len(df) > MAX_ROWS
                if truncated_sheet:
                    df = df.head(MAX_ROWS)

                result[sheet] = {
                    "columns":   df.columns.tolist(),
                    "rows":      df.values.tolist(),
                    "shape":     [len(df), len(df.columns)],
                    "truncated": truncated_sheet,
                    "dtypes":    {c: str(t) for c, t in df.dtypes.items()},
                }

            # Summary stats for numeric columns in first sheet
            first_df  = xl.parse(sheets[0]) if sheets else None
            stats     = {}
            if first_df is not None:
                num_cols = first_df.select_dtypes(include="number").columns.tolist()
                if num_cols:
                    stats = first_df[num_cols].describe().round(2).to_dict()

            return {
                "status":      "success",
                "path":        str(p),
                "sheets":      sheets,
                "sheet_count": len(sheets),
                "data":        result,
                "stats":       stats,
            }
        except ImportError:
            return {"status": "error", "error": "pandas not installed. Run: pip install pandas openpyxl"}
        except Exception as e:
            return {"status": "error", "error": f"XLSX read failed: {type(e).__name__}: {e}"}

    # ── write_xlsx ────────────────────────────────────────────────────────────
    if action == "write_xlsx":
        p, err = _safe_resolve(path)
        if err:
            return {"status": "error", "error": err}
        if not content and not isinstance(content, (dict, list)):
            return {"status": "error", "error": "content is required for write_xlsx"}
        if p.suffix.lower() not in (".xlsx", ".xls"):
            p = p.with_suffix(".xlsx")

        try:
            import pandas as pd
            from openpyxl.styles import Font, PatternFill, Alignment
            from openpyxl.utils import get_column_letter

            # content can be:
            # A) {"Sheet1": [{"col": val, ...}, ...], "Sheet2": [...]}  ← multi-sheet
            # B) [{"col": val, ...}, ...]                               ← single sheet
            # C) {"columns": [...], "rows": [[...], ...]}               ← raw format
            # D) a plain string (JSON) — parse it
            if isinstance(content, str):
                import json as _json
                try:
                    content = _json.loads(content)
                except Exception:
                    return {"status": "error",
                            "error": "content string could not be parsed as JSON"}

            p.parent.mkdir(parents=True, exist_ok=True)

            with pd.ExcelWriter(str(p), engine="openpyxl") as writer:
                sheets_written = []

                if isinstance(content, dict):
                    # Normalise keys to strings to prevent DataFrame errors
                    content = {str(k): v for k, v in content.items()}

                    # Check if it's multi-sheet (values are lists) or raw format
                    first_val = next(iter(content.values()), None)
                    if isinstance(first_val, list) and all(
                        isinstance(v, list) for v in content.values()
                    ):
                        # Multi-sheet: {"Sheet1": [rows...], "Sheet2": [rows...]}
                        for sheet_name, rows in content.items():
                            safe_name = str(sheet_name)[:31]
                            try:
                                if rows and isinstance(rows[0], dict):
                                    # Normalise row keys too
                                    rows = [{str(k): v for k, v in r.items()} for r in rows]
                                    df = pd.DataFrame(rows)
                                elif rows and isinstance(rows[0], list):
                                    df = pd.DataFrame(rows[1:], columns=rows[0])
                                else:
                                    df = pd.DataFrame(rows)
                                df.to_excel(writer, sheet_name=safe_name, index=False)
                                sheets_written.append(sheet_name)
                            except Exception as sheet_err:
                                # Skip bad sheet, continue with others
                                sheets_written.append(f"{safe_name}(error:{sheet_err})")
                    elif "columns" in content and "rows" in content:
                        # Raw format
                        df = pd.DataFrame(content["rows"],
                                          columns=content["columns"])
                        sheet_name = str(content.get("sheet", "Sheet1"))[:31]
                        df.to_excel(writer, sheet_name=sheet_name, index=False)
                        sheets_written.append(sheet_name)
                    else:
                        # Single sheet from dict of lists -- normalise keys
                        safe_content = {str(k): v for k, v in content.items()}
                        df = pd.DataFrame(safe_content)
                        df.to_excel(writer, sheet_name="Sheet1", index=False)
                        sheets_written.append("Sheet1")

                elif isinstance(content, list):
                    if content and isinstance(content[0], dict):
                        df = pd.DataFrame(content)
                    elif content and isinstance(content[0], list):
                        df = pd.DataFrame(content[1:], columns=content[0])
                    else:
                        df = pd.DataFrame(content)
                    df.to_excel(writer, sheet_name="Sheet1", index=False)
                    sheets_written.append("Sheet1")

                # Style header rows
                wb = writer.book
                for ws in wb.worksheets:
                    for cell in ws[1]:  # first row = header
                        cell.font      = Font(bold=True, color="FFFFFF")
                        cell.fill      = PatternFill("solid", fgColor="2C3E50")
                        cell.alignment = Alignment(horizontal="center")
                    # Auto-width
                    for col_idx, col in enumerate(ws.columns, 1):
                        max_len = max(
                            (len(str(c.value)) for c in col if c.value is not None),
                            default=8,
                        )
                        ws.column_dimensions[get_column_letter(col_idx)].width = min(
                            max_len + 4, 40
                        )

            return {
                "status":         "success",
                "path":           str(p),
                "size":           p.stat().st_size,
                "sheets_written": sheets_written,
            }
        except ImportError:
            return {"status": "error",
                    "error": "pandas/openpyxl not installed. Run: pip install pandas openpyxl"}
        except Exception as e:
            return {"status": "error", "error": f"XLSX write failed: {type(e).__name__}: {e}"}

    # ── read_pptx ─────────────────────────────────────────────────────────────
    if action == "read_pptx":
        p, err = _safe_resolve(path)
        if err:
            return {"status": "error", "error": err}
        if not p.exists():
            return {"status": "error", "error": f"File not found: {p}"}
        if p.suffix.lower() != ".pptx":
            return {"status": "error", "error": f"Not a .pptx file: {p.name}"}
        try:
            from pptx import Presentation
            from pptx.util import Pt

            prs    = Presentation(str(p))
            slides = []

            for i, slide in enumerate(prs.slides, 1):
                texts  = []
                images = 0
                tables = 0

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                texts.append(text)
                    if hasattr(shape, "table"):
                        tables += 1
                        # Extract table text
                        for row in shape.table.rows:
                            cells = [c.text.strip() for c in row.cells]
                            texts.append(" | ".join(cells))
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        images += 1

                # Slide layout name
                layout = ""
                try:
                    layout = slide.slide_layout.name
                except Exception:
                    pass

                slides.append({
                    "slide":  i,
                    "layout": layout,
                    "texts":  texts,
                    "images": images,
                    "tables": tables,
                })

            # Flat readable text
            flat = "\n\n".join(
                f"--- Slide {s['slide']} ---\n" + "\n".join(s["texts"])
                for s in slides
                if s["texts"]
            )

            truncated = len(flat) > max_chars
            if truncated:
                flat = flat[:max_chars] + "\n\n[...truncated]"

            return {
                "status":      "success",
                "path":        str(p),
                "text":        flat,
                "slides":      slides,
                "slide_count": len(slides),
                "truncated":   truncated,
            }
        except ImportError:
            return {"status": "error",
                    "error": "python-pptx not installed. Run: pip install python-pptx"}
        except Exception as e:
            return {"status": "error", "error": f"PPTX read failed: {type(e).__name__}: {e}"}

    # ── write_pptx ────────────────────────────────────────────────────────────
    if action == "write_pptx":
        p, err = _safe_resolve(path)
        if err:
            return {"status": "error", "error": err}
        if not content:
            return {"status": "error", "error": "content is required for write_pptx"}
        if p.suffix.lower() != ".pptx":
            p = p.with_suffix(".pptx")

        # content = list of slide dicts:
        # {"title": "...", "body": "...", "bullets": [...], "layout": "title|content|blank"}
        # OR a JSON string
        if isinstance(content, str):
            import json as _json
            try:
                content = _json.loads(content)
            except Exception:
                return {"status": "error",
                        "error": "content string could not be parsed as JSON"}

        if not isinstance(content, list):
            return {"status": "error",
                    "error": "content must be a list of slide dicts: [{title, body/bullets}, ...]"}

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            prs = Presentation()

            # Slide dimensions (16:9 widescreen)
            prs.slide_width  = Inches(13.33)
            prs.slide_height = Inches(7.5)

            # Color scheme
            DARK   = RGBColor(0x2C, 0x3E, 0x50)  # dark blue-grey
            ACCENT = RGBColor(0x34, 0x98, 0xDB)  # blue
            WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
            LIGHT  = RGBColor(0xEC, 0xF0, 0xF1)

            layouts = {n.name: n for n in prs.slide_layouts}

            # ── helpers ───────────────────────────────────────────────────────
            def _add_rect(slide, left, top, width, height, color):
                from pptx.util import Emu
                shape = slide.shapes.add_shape(
                    1,  # MSO_SHAPE_TYPE.RECTANGLE
                    Inches(left), Inches(top), Inches(width), Inches(height),
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = color
                shape.line.fill.background()
                return shape

            def _add_text(slide, text, left, top, width, height,
                          font_size=18, bold=False, color=None, align="left"):
                txb  = slide.shapes.add_textbox(
                    Inches(left), Inches(top), Inches(width), Inches(height))
                tf   = txb.text_frame
                tf.word_wrap = True
                para = tf.paragraphs[0]
                para.alignment = {
                    "left":   PP_ALIGN.LEFT,
                    "center": PP_ALIGN.CENTER,
                    "right":  PP_ALIGN.RIGHT,
                }.get(align, PP_ALIGN.LEFT)
                run           = para.add_run()
                run.text      = text
                run.font.size = Pt(font_size)
                run.font.bold = bold
                if color:
                    run.font.color.rgb = color
                return txb

            # ── build slides ──────────────────────────────────────────────────
            for i, spec in enumerate(content):
                slide_title   = spec.get("title", f"Slide {i+1}")
                slide_body    = spec.get("body", "")
                slide_bullets = spec.get("bullets", [])
                slide_layout  = spec.get("layout", "content").lower()
                slide_notes   = spec.get("notes", "")
                is_title_slide = (slide_layout == "title" or i == 0 and slide_layout != "content")

                # Use blank layout for full control
                blank = prs.slide_layouts[6]  # blank
                slide = prs.slides.add_slide(blank)

                if is_title_slide and i == 0:
                    # ── Title slide — dark background, centered ──────────────
                    _add_rect(slide, 0, 0, 13.33, 7.5, DARK)
                    _add_rect(slide, 0, 3.0, 13.33, 0.06, ACCENT)
                    _add_text(slide, slide_title,
                              1, 2.2, 11.33, 1.4,
                              font_size=40, bold=True, color=WHITE, align="center")
                    if slide_body:
                        _add_text(slide, slide_body,
                                  1, 3.8, 11.33, 0.8,
                                  font_size=20, color=LIGHT, align="center")
                    if spec.get("subtitle"):
                        _add_text(slide, spec["subtitle"],
                                  1, 4.6, 11.33, 0.7,
                                  font_size=16, color=ACCENT, align="center")
                else:
                    # ── Content slide ─────────────────────────────────────────
                    # Top accent bar
                    _add_rect(slide, 0, 0, 13.33, 1.1, DARK)
                    _add_rect(slide, 0, 1.1, 13.33, 0.05, ACCENT)

                    # Title
                    _add_text(slide, slide_title,
                              0.3, 0.15, 12.73, 0.85,
                              font_size=26, bold=True, color=WHITE)

                    # Slide number
                    _add_text(slide, str(i + 1),
                              12.5, 0.25, 0.6, 0.5,
                              font_size=13, color=ACCENT, align="right")

                    content_top = 1.35

                    if slide_bullets:
                        # Bullet list
                        txb = slide.shapes.add_textbox(
                            Inches(0.5), Inches(content_top),
                            Inches(12.33), Inches(5.8))
                        tf  = txb.text_frame
                        tf.word_wrap = True

                        for j, bullet in enumerate(slide_bullets):
                            if j == 0:
                                para = tf.paragraphs[0]
                            else:
                                para = tf.add_paragraph()

                            # Support nested bullets: {"text":"...", "level":1}
                            if isinstance(bullet, dict):
                                bullet_text  = bullet.get("text", "")
                                bullet_level = bullet.get("level", 0)
                            else:
                                bullet_text  = str(bullet)
                                bullet_level = 0

                            para.level = bullet_level
                            run = para.add_run()
                            run.text      = bullet_text
                            run.font.size = Pt(18 - bullet_level * 2)
                            run.font.color.rgb = DARK
                            para.space_before = Pt(6)

                    elif slide_body:
                        _add_text(slide, slide_body,
                                  0.5, content_top, 12.33, 5.8,
                                  font_size=18, color=DARK)

                # Speaker notes
                if slide_notes:
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = slide_notes

            p.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(p))

            return {
                "status":      "success",
                "path":        str(p),
                "size":        p.stat().st_size,
                "slide_count": len(content),
            }
        except ImportError:
            return {"status": "error",
                    "error": "python-pptx not installed. Run: pip install python-pptx"}
        except Exception as e:
            return {"status": "error",
                    "error": f"PPTX write failed: {type(e).__name__}: {e}"}

    # ── write_pdf ─────────────────────────────────────────────────────────────
    if action == "write_pdf":
        p, err = _safe_resolve(path)
        if err:
            return {"status": "error", "error": err}
        if not content:
            return {"status": "error", "error": "content is required for write_pdf"}
        if p.suffix.lower() != ".pdf":
            p = p.with_suffix(".pdf")

        try:
            from fpdf import FPDF

            pdf = FPDF()
            pdf.set_auto_page_break(auto=True, margin=15)
            pdf.add_page()

            if title:
                pdf.set_title(title)
                pdf.set_font("Helvetica", "B", 18)
                pdf.cell(0, 12, title, ln=True)
                pdf.ln(4)

            pdf.set_font("Helvetica", size=11)

            for line in content.splitlines():
                stripped = line.strip()
                if stripped.startswith("## "):
                    pdf.set_font("Helvetica", "B", 14)
                    pdf.ln(3)
                    pdf.cell(0, 9, stripped[3:], ln=True)
                    pdf.set_font("Helvetica", size=11)
                elif stripped.startswith("# "):
                    pdf.set_font("Helvetica", "B", 16)
                    pdf.ln(4)
                    pdf.cell(0, 10, stripped[2:], ln=True)
                    pdf.set_font("Helvetica", size=11)
                elif stripped.startswith("### "):
                    pdf.set_font("Helvetica", "B", 12)
                    pdf.ln(2)
                    pdf.cell(0, 8, stripped[4:], ln=True)
                    pdf.set_font("Helvetica", size=11)
                elif stripped in ("---", "***"):
                    pdf.ln(2)
                    pdf.line(pdf.l_margin, pdf.get_y(),
                             pdf.w - pdf.r_margin, pdf.get_y())
                    pdf.ln(2)
                elif not stripped:
                    pdf.ln(4)
                else:
                    pdf.multi_cell(0, 6, line)

            p.parent.mkdir(parents=True, exist_ok=True)
            pdf.output(str(p))
            return {
                "status": "success",
                "path":   str(p),
                "pages":  pdf.page,
                "size":   p.stat().st_size,
            }
        except ImportError:
            return {"status": "error", "error": "fpdf2 not installed. Run: pip install fpdf2"}
        except Exception as e:
            return {"status": "error", "error": f"PDF write failed: {type(e).__name__}: {e}"}

    # ── write_docx ────────────────────────────────────────────────────────────
    if action == "write_docx":
        p, err = _safe_resolve(path)
        if err:
            return {"status": "error", "error": err}
        if not content:
            return {"status": "error", "error": "content is required for write_docx"}
        if p.suffix.lower() != ".docx":
            p = p.with_suffix(".docx")

        try:
            import re as _re
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            doc = Document()

            if title:
                t = doc.add_heading(title, level=0)
                t.alignment = WD_ALIGN_PARAGRAPH.CENTER

            for line in content.splitlines():
                stripped = line.strip()
                if stripped.startswith("### "):
                    doc.add_heading(stripped[4:], level=3)
                elif stripped.startswith("## "):
                    doc.add_heading(stripped[3:], level=2)
                elif stripped.startswith("# "):
                    doc.add_heading(stripped[2:], level=1)
                elif stripped.startswith(("- ", "* ")):
                    doc.add_paragraph(stripped[2:], style="List Bullet")
                elif _re.match(r"^\d+\.\s", stripped):
                    doc.add_paragraph(_re.sub(r"^\d+\.\s", "", stripped),
                                      style="List Number")
                elif stripped in ("---", "***"):
                    doc.add_paragraph("_" * 60)
                elif not stripped:
                    doc.add_paragraph("")
                else:
                    para  = doc.add_paragraph()
                    parts = _re.split(r"(\*\*[^*]+\*\*)", stripped)
                    for part in parts:
                        if part.startswith("**") and part.endswith("**"):
                            para.add_run(part[2:-2]).bold = True
                        else:
                            para.add_run(part)

            p.parent.mkdir(parents=True, exist_ok=True)
            doc.save(str(p))
            return {
                "status":     "success",
                "path":       str(p),
                "size":       p.stat().st_size,
                "paragraphs": len(doc.paragraphs),
            }
        except ImportError:
            return {"status": "error",
                    "error": "python-docx not installed. Run: pip install python-docx"}
        except Exception as e:
            return {"status": "error",
                    "error": f"DOCX write failed: {type(e).__name__}: {e}"}

    # ── compress ──────────────────────────────────────────────────────────────
    if action == "compress":
        import zipfile
        p, err = _safe_resolve(path or ".")
        if err:
            return {"status": "error", "error": err}
        if not p.exists():
            return {"status": "error", "error": f"Path not found: {p}"}

        # Output zip sits next to the target (or in workspace root if target is root)
        zip_name = (p.name or "workspace") + ".zip"
        zip_path = (p.parent if p.parent != p else cfg.workspace_root) / zip_name

        try:
            file_count = 0
            with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
                if p.is_dir():
                    for item in p.rglob("*"):
                        if item.is_file():
                            arcname = item.relative_to(p)
                            zf.write(item, arcname)
                            file_count += 1
                else:
                    zf.write(p, p.name)
                    file_count = 1

            return {
                "status":     "success",
                "zip_path":   str(zip_path),
                "files":      file_count,
                "size":       zip_path.stat().st_size,
            }
        except Exception as e:
            return {"status": "error", "error": f"Compress failed: {e}"}

    return {
        "status": "error",
        "error":  (
            f"Unknown action '{action}'. "
            "Use: read | write | list | backup | read_many | search | "
            "read_pdf | write_pdf | "
            "read_docx | write_docx | "
            "read_xlsx | write_xlsx | "
            "read_pptx | write_pptx | compress"
        ),
    }
